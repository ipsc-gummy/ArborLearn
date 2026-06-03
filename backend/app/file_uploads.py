from __future__ import annotations

import base64
import json
import re
import socket
import time
from io import BytesIO
from pathlib import Path
from urllib import error, request

from fastapi import HTTPException, UploadFile

from .settings import (
    get_max_upload_bytes,
    get_ocr_languages,
    get_ocr_timeout_seconds,
    get_tesseract_cmd,
    get_upload_dir,
    get_vision_api_key,
    get_vision_base_url,
    get_vision_max_attempts,
    get_vision_max_image_edge,
    get_vision_model,
    get_vision_provider,
    get_vision_retry_initial_delay_seconds,
    get_vision_retry_max_delay_seconds,
    get_vision_timeout_seconds,
    is_ocr_enabled,
)


ALLOWED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
}
MAX_EXTRACTED_CHARS = 120_000
DEFERRED_EXTRACTION_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg", ".webp", ".bmp"}


class VisionRequestError(RuntimeError):
    def __init__(self, message: str, *, retryable: bool = True):
        super().__init__(message)
        self.retryable = retryable


def safe_filename(filename: str | None) -> str:
    raw_name = Path(filename or "upload.txt").name.strip()
    if not raw_name:
        raw_name = "upload.txt"
    stem = Path(raw_name).stem
    suffix = Path(raw_name).suffix.lower()
    clean_stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._-")
    if not clean_stem:
        clean_stem = "upload"
    return f"{clean_stem[:80]}{suffix}"


def validate_upload_name(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_EXTENSIONS))
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {allowed}")
    return suffix


def should_defer_extraction(suffix: str) -> bool:
    return suffix in DEFERRED_EXTRACTION_EXTENSIONS


def _trim(text: str) -> str:
    return text[:MAX_EXTRACTED_CHARS]


def _json_payload(payload: dict) -> str:
    return _trim(json.dumps(payload, ensure_ascii=False, indent=2))


def _vision_retry_delay_seconds(attempt: int) -> float:
    initial_delay = get_vision_retry_initial_delay_seconds()
    max_delay = get_vision_retry_max_delay_seconds()
    if initial_delay <= 0 or max_delay <= 0:
        return 0.0
    return min(max_delay, initial_delay * (2 ** max(0, attempt - 1)))


def _image_resampling_filter():
    try:
        from PIL import Image

        return Image.Resampling.LANCZOS
    except Exception:
        from PIL import Image

        return Image.LANCZOS


def _prepare_vision_image_bytes(content: bytes, image, mime_type: str) -> tuple[bytes, str, dict, list[str]]:
    width, height = image.size
    max_edge = get_vision_max_image_edge()
    metadata = {
        "original_size": {"width": width, "height": height},
        "input_size": {"width": width, "height": height},
        "preprocessed": False,
        "max_image_edge": max_edge,
    }
    warnings: list[str] = []
    if not max_edge or max(width, height) <= max_edge:
        return content, mime_type, metadata, warnings

    resized = image.copy()
    if resized.mode not in {"RGB", "RGBA", "L"}:
        resized = resized.convert("RGB")
    resized.thumbnail((max_edge, max_edge), _image_resampling_filter())
    output = BytesIO()
    output_format = "PNG"
    resized.save(output, format=output_format, optimize=True)
    metadata.update(
        {
            "input_size": {"width": resized.width, "height": resized.height},
            "preprocessed": True,
            "preprocess_reason": f"max edge {max(width, height)}px exceeds {max_edge}px",
        }
    )
    warnings.append(
        f"Vision input resized from {width}x{height} to {resized.width}x{resized.height}"
    )
    return output.getvalue(), "image/png", metadata, warnings


def decode_text_file(content: bytes) -> tuple[str, str, str | None]:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            text = content.decode(encoding)
            payload = {
                "type": "text",
                "text": _trim(text),
                "tables": [],
                "images_summary": "",
                "metadata": {"encoding": encoding},
            }
            return _json_payload(payload), "ready", None
        except UnicodeDecodeError:
            continue
    return "", "failed", "Unable to decode file as text"


def decode_pdf_file(content: bytes, filename: str) -> tuple[str, str, str | None]:
    extracted = ""
    parser_used = "none"
    tables: list[dict] = []
    image_notes: list[str] = []
    warnings: list[str] = []

    try:
        import fitz  # PyMuPDF

        parser_used = "pymupdf"
        doc = fitz.open(stream=content, filetype="pdf")
        parts: list[str] = []
        for idx, page in enumerate(doc):
            page_text = (page.get_text("text") or "").strip()
            if page_text:
                parts.append(page_text)
            image_count = len(page.get_images(full=True))
            if image_count:
                image_notes.append(f"page {idx+1}: {image_count} embedded images")
        extracted = "\n\n".join(parts).strip()
    except Exception as exc:
        warnings.append(f"PyMuPDF unavailable/failed: {exc}")

    if not extracted:
        try:
            from pypdf import PdfReader

            parser_used = "pypdf"
            reader = PdfReader(BytesIO(content))
            parts = [(page.extract_text() or "").strip() for page in reader.pages]
            extracted = "\n\n".join(part for part in parts if part).strip()
        except Exception as exc:
            warnings.append(f"pypdf fallback failed: {exc}")

    if not extracted:
        return "", "failed", "No extractable text found in PDF"

    payload = {
        "type": "pdf",
        "text": _trim(extracted),
        "tables": tables,
        "images_summary": "; ".join(image_notes),
        "metadata": {"filename": filename, "parser": parser_used, "warnings": warnings},
    }
    return _json_payload(payload), "ready", None


def decode_docx_file(content: bytes, filename: str) -> tuple[str, str, str | None]:
    try:
        from docx import Document

        document = Document(BytesIO(content))
        text_parts: list[str] = [(p.text or "").strip() for p in document.paragraphs if (p.text or "").strip()]
        tables: list[dict] = []
        for table in document.tables:
            rows = []
            for row in table.rows:
                row_values = [(cell.text or "").strip() for cell in row.cells]
                if any(row_values):
                    rows.append(row_values)
            if rows:
                tables.append({"rows": rows})
        payload = {
            "type": "docx",
            "text": _trim("\n".join(text_parts)),
            "tables": tables,
            "images_summary": "",
            "metadata": {"filename": filename, "paragraph_count": len(text_parts), "table_count": len(tables)},
        }
        return _json_payload(payload), "ready", None
    except Exception as exc:
        return "", "failed", f"Unable to parse DOCX: {exc}"


def decode_xlsx_file(content: bytes, filename: str) -> tuple[str, str, str | None]:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(BytesIO(content), data_only=True, read_only=True)
        lines: list[str] = []
        tables: list[dict] = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell).strip() for cell in row]
                if any(values):
                    rows.append(values)
            if rows:
                header = rows[0]
                body = rows[1:101]
                tables.append({"sheet": ws.title, "header": header, "rows": body})
                lines.append(f"[Sheet] {ws.title}")
                lines.extend(" | ".join(r) for r in rows[:30])
        payload = {
            "type": "xlsx",
            "text": _trim("\n".join(lines)),
            "tables": tables,
            "images_summary": "",
            "metadata": {"filename": filename, "sheet_count": len(wb.worksheets)},
        }
        return _json_payload(payload), "ready", None
    except Exception as exc:
        return "", "failed", f"Unable to parse XLSX: {exc}"


def decode_pptx_file(content: bytes, filename: str) -> tuple[str, str, str | None]:
    try:
        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        slide_summaries = []
        text_parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            merged = "\n".join(t for t in slide_text if t)
            if merged:
                text_parts.append(f"[Slide {i}]\n{merged}")
            slide_summaries.append({"slide": i, "text_blocks": len(slide_text)})
        payload = {
            "type": "pptx",
            "text": _trim("\n\n".join(text_parts)),
            "tables": [],
            "images_summary": "",
            "metadata": {"filename": filename, "slide_count": len(slide_summaries), "slides": slide_summaries},
        }
        return _json_payload(payload), "ready", None
    except Exception as exc:
        return "", "failed", f"Unable to parse PPTX: {exc}"


def _call_vision_model(image_bytes: bytes, mime_type: str) -> str:
    provider = get_vision_provider()
    if provider in {"", "none", "disabled"}:
        return ""
    if provider not in {"openai_compatible", "qwen_vl"}:
        raise VisionRequestError(f"Unsupported VISION_PROVIDER: {provider}", retryable=False)

    b64 = base64.b64encode(image_bytes).decode("ascii")
    data_url = f"data:{mime_type};base64,{b64}"
    payload = {
        "model": get_vision_model(),
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "请描述这张图片，提取其中关键实体、数字和上下文，返回中文结构化摘要。"},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ],
        "temperature": 0.1,
    }
    body = json.dumps(payload).encode("utf-8")
    api_key = get_vision_api_key()
    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    req = request.Request(
        url=f"{get_vision_base_url()}/chat/completions",
        data=body,
        headers=headers,
        method="POST",
    )
    try:
        with request.urlopen(req, timeout=get_vision_timeout_seconds()) as resp:
            result = json.loads(resp.read().decode("utf-8"))
    except error.HTTPError as exc:
        retryable = exc.code >= 500 or exc.code in {408, 409, 425, 429}
        detail = ""
        try:
            detail = exc.read(512).decode("utf-8", errors="replace").strip()
        except Exception:
            detail = ""
        suffix = f": {detail}" if detail else ""
        raise VisionRequestError(
            f"Vision request failed with HTTP {exc.code}{suffix}",
            retryable=retryable,
        ) from exc
    except error.URLError as exc:
        raise VisionRequestError(f"Vision request failed: {exc}", retryable=True) from exc
    except (TimeoutError, socket.timeout) as exc:
        raise VisionRequestError(f"Vision decode failed: {exc}", retryable=True) from exc
    except json.JSONDecodeError as exc:
        raise VisionRequestError(f"Vision decode failed: {exc}", retryable=True) from exc
    except Exception as exc:
        raise VisionRequestError(f"Vision decode failed: {exc}", retryable=True) from exc

    try:
        return (result["choices"][0]["message"]["content"] or "").strip()
    except Exception as exc:
        raise VisionRequestError(f"Vision response format error: {exc}", retryable=False) from exc


def decode_image_file(content: bytes, filename: str, mime_type: str | None) -> tuple[str, str, str | None]:
    try:
        from PIL import Image

        image = Image.open(BytesIO(content))
        image.load()
        width, height = image.size
        mime = mime_type or f"image/{(image.format or 'png').lower()}"
        ocr_text = ""
        vision_summary = ""
        vision_error: str | None = None
        warnings: list[str] = []
        vision_input_bytes = content
        vision_mime = mime
        vision_metadata = {
            "original_size": {"width": width, "height": height},
            "input_size": {"width": width, "height": height},
            "preprocessed": False,
            "max_image_edge": get_vision_max_image_edge(),
        }

        # OCR (optional fallback channel)
        if is_ocr_enabled():
            try:
                import pytesseract

                tesseract_cmd = get_tesseract_cmd()
                if tesseract_cmd:
                    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
                ocr_text = (
                    pytesseract.image_to_string(
                        image,
                        lang=get_ocr_languages(),
                        timeout=get_ocr_timeout_seconds(),
                    )
                    or ""
                ).strip()
            except Exception as exc:
                warnings.append(f"OCR unavailable: {exc}")
        else:
            warnings.append("OCR disabled by configuration")

        # Vision model (primary for image understanding)
        provider = get_vision_provider()
        if provider in {"", "none", "disabled"}:
            warnings.append("Vision provider disabled")
        else:
            vision_input_bytes, vision_mime, vision_metadata, preprocess_warnings = _prepare_vision_image_bytes(
                content,
                image,
                mime,
            )
            warnings.extend(preprocess_warnings)
            max_attempts = get_vision_max_attempts()
            for attempt in range(1, max_attempts + 1):
                started_at = time.perf_counter()
                retryable = True
                try:
                    vision_summary = _call_vision_model(vision_input_bytes, vision_mime)
                    elapsed_ms = int((time.perf_counter() - started_at) * 1000)
                    if vision_summary:
                        vision_error = None
                        break
                    vision_error = "Vision provider returned empty summary"
                except Exception as exc:
                    elapsed_ms = int((time.perf_counter() - started_at) * 1000)
                    vision_error = str(exc)
                    retryable = getattr(exc, "retryable", True)
                if attempt == max_attempts or not retryable:
                    warnings.append(f"Vision attempt {attempt} failed after {elapsed_ms}ms: {vision_error}")
                    break
                delay = _vision_retry_delay_seconds(attempt)
                warnings.append(
                    f"Vision attempt {attempt} failed after {elapsed_ms}ms: {vision_error}; retrying in {delay:g}s"
                )
                if delay > 0:
                    time.sleep(delay)

        payload = {
            "type": "image",
            "text": _trim(ocr_text),
            "tables": [],
            "images_summary": _trim(vision_summary),
            "metadata": {
                "filename": filename,
                "format": image.format or "unknown",
                "size": {"width": width, "height": height},
                "mime_type": mime,
                "vision_input": {**vision_metadata, "mime_type": vision_mime},
                "vision_provider": get_vision_provider(),
                "vision_max_attempts": get_vision_max_attempts(),
                "vision_timeout_seconds": get_vision_timeout_seconds(),
                "warnings": warnings,
            },
        }
        has_vision_signal = bool(vision_summary.strip())
        if get_vision_provider() not in {"", "none", "disabled"} and vision_error and not has_vision_signal:
            return _json_payload(payload), "failed", vision_error
        return _json_payload(payload), "ready", None
    except Exception as exc:
        return "", "failed", f"Unable to parse image: {exc}"


def extract_file_text(content: bytes, suffix: str, filename: str, mime_type: str | None) -> tuple[str, str, str | None]:
    if suffix in {".txt", ".md"}:
        return decode_text_file(content)
    if suffix == ".pdf":
        return decode_pdf_file(content, filename)
    if suffix == ".docx":
        return decode_docx_file(content, filename)
    if suffix == ".xlsx":
        return decode_xlsx_file(content, filename)
    if suffix == ".pptx":
        return decode_pptx_file(content, filename)
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}:
        return decode_image_file(content, filename, mime_type)
    return "", "failed", "Unsupported file type"


def extract_stored_file(*, storage_path: str, filename: str, mime_type: str | None) -> tuple[str, str, str | None]:
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        return "", "failed", "Unsupported file type"

    path = Path(storage_path)
    if not path.exists():
        return "", "failed", "Stored file not found"

    content = path.read_bytes()
    if not content:
        return "", "failed", "Stored file is empty"

    return extract_file_text(content, suffix, filename, mime_type)


async def prepare_uploaded_file(upload: UploadFile, *, user_id: str, file_id: str) -> dict:
    filename = safe_filename(upload.filename)
    suffix = validate_upload_name(filename)

    content = await upload.read()
    max_upload_bytes = get_max_upload_bytes()
    if len(content) > max_upload_bytes:
        max_mb = max_upload_bytes // (1024 * 1024)
        raise HTTPException(status_code=400, detail=f"File is too large. Max upload size is {max_mb} MB")
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    user_dir = get_upload_dir() / user_id / file_id
    user_dir.mkdir(parents=True, exist_ok=True)
    storage_path = user_dir / filename
    storage_path.write_bytes(content)

    if should_defer_extraction(suffix):
        extracted_text = ""
        extraction_status = "pending"
        error_message = None
    else:
        extracted_text, extraction_status, error_message = extract_file_text(content, suffix, filename, upload.content_type)

    return {
        "filename": filename,
        "original_filename": upload.filename or filename,
        "mime_type": upload.content_type,
        "file_size": len(content),
        "storage_path": str(storage_path),
        "extracted_text": extracted_text,
        "extraction_status": extraction_status,
        "error_message": error_message,
    }
